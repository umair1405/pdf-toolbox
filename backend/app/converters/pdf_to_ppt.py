import fitz
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def pdf_to_ppt(pdf_path, output_dir):
    pdf = fitz.open(pdf_path)

    prs = Presentation()

    blank = prs.slide_layouts[6]

    for page in pdf:
        slide = prs.slides.add_slide(blank)

        blocks = page.get_text("dict")["blocks"]

        for block in blocks:
            if "lines" not in block:
                continue

            x0, y0, x1, y1 = block["bbox"]

            textbox = slide.shapes.add_textbox(
                Pt(x0),
                Pt(y0),
                Pt(max(x1 - x0, 100)),
                Pt(max(y1 - y0, 30)),
            )

            tf = textbox.text_frame

            for line in block["lines"]:
                text = ""

                for span in line["spans"]:
                    text += span["text"]

                if text.strip():
                    p = tf.add_paragraph()
                    p.text = text
                    p.font.size = Pt(18)

    output = (
        Path(output_dir)
        / f"{Path(pdf_path).stem}.pptx"
    )

    prs.save(str(output))

    return str(output)